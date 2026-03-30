from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Paleta AuroraISP
AZUL = RGBColor(0x25, 0x63, 0xEB)       # #2563EB
DARK = RGBColor(0x0F, 0x17, 0x2A)       # #0F172A
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF1, 0xF5, 0xF9)
CINZA_TEXTO = RGBColor(0x64, 0x74, 0x8B)
CIANO = RGBColor(0x06, 0xB6, 0xD4)      # acento gradiente

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout em branco


def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_paragraph(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


# ──────────────────────────────────────────────
# SLIDE 1 — CAPA
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK)
add_rect(slide, 0, 5.8, 13.33, 1.7, AZUL)

add_text(slide, "AuroraISP", 1, 1.2, 11, 1.2, 52, BRANCO, bold=True)
add_text(slide, "Vende mais. Perde menos. Fideliza sempre.", 1, 2.5, 11, 0.8, 22, CIANO, bold=False)
add_text(slide, "Hub de tecnologia com IA para provedores de internet.", 1, 3.4, 10, 0.7, 16, CINZA_CLARO)
add_text(slide, "Pitch Comercial — Confidencial", 1, 6.1, 11, 0.6, 12, BRANCO)


# ──────────────────────────────────────────────
# SLIDE 2 — O PROBLEMA
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, CINZA_CLARO)
add_rect(slide, 0, 0, 13.33, 1.1, DARK)
add_text(slide, "O dia a dia do seu time comercial hoje", 0.5, 0.2, 12, 0.8, 22, BRANCO, bold=True)

problemas = [
    ("Lead some no WhatsApp",        "Chega, pergunta o preço e desaparece. Ninguém acompanha."),
    ("Documentos no caos",           "Vendedor pede CPF no chat, o cliente manda foto borrada, o processo trava."),
    ("Ativação manual e lenta",      "Depois do contrato, alguém precisa lançar tudo no sistema. Retrabalho todo dia."),
    ("Time grande para pouco volume","8 pessoas para atender 1.000 leads por mês. Custo alto, resultado médio."),
]

x_positions = [0.4, 3.6, 6.8, 10.0]
for i, (titulo, desc) in enumerate(problemas):
    x = x_positions[i]
    add_rect(slide, x, 1.4, 3.0, 4.8, BRANCO)
    add_rect(slide, x, 1.4, 3.0, 0.6, AZUL)
    add_text(slide, titulo, x + 0.15, 1.5, 2.7, 0.5, 13, BRANCO, bold=True)
    add_text(slide, desc, x + 0.15, 2.2, 2.7, 3.5, 12, CINZA_TEXTO)


# ──────────────────────────────────────────────
# SLIDE 3 — A SOLUÇÃO
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK)
add_rect(slide, 0, 0, 13.33, 1.1, AZUL)
add_text(slide, "A AuroraISP resolve isso", 0.5, 0.2, 12, 0.8, 22, BRANCO, bold=True)

add_text(slide, "Do lead ao cliente fidelizado com inteligência.", 0.8, 1.4, 11, 0.8, 20, CIANO, bold=True)
add_text(slide,
    "A Aurora automatiza o fluxo comercial do início ao fim. O lead entra no WhatsApp, o bot qualifica, "
    "coleta os documentos, gera o contrato e ativa direto no HubSoft. Sem planilha, sem retrabalho, sem lead perdido.",
    0.8, 2.3, 11.5, 1.5, 15, CINZA_CLARO)

add_text(slide, "Integração nativa e profunda com HubSoft.", 0.8, 4.0, 11, 0.7, 15, BRANCO, bold=True)
add_text(slide, "Nenhum concorrente faz isso.", 0.8, 4.6, 11, 0.6, 13, CINZA_CLARO, italic=True)


# ──────────────────────────────────────────────
# SLIDE 4 — OS 3 MÓDULOS
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, CINZA_CLARO)
add_rect(slide, 0, 0, 13.33, 1.1, DARK)
add_text(slide, "Um hub. Três módulos. Uma jornada completa.", 0.5, 0.2, 12, 0.8, 22, BRANCO, bold=True)

modulos = [
    ("Comercial", "✅ Disponível",
     "Lead no WhatsApp\nBot de qualificação\nColeta de documentos\nGeração de contrato\nAtivação no HubSoft\nCRM Kanban (Pro)"),
    ("Marketing", "🔧 Em breve",
     "Disparos de WhatsApp e e-mail\nFluxos de automação prontos\nNutrição de leads\nGestão de tráfego pago com IA\nRelatórios de atribuição"),
    ("CS", "🔧 Em breve",
     "Clube de Benefícios\nPrevenção de churn\nNPS automatizado\nUpsell de plano e SVA\nRelatórios de saúde da base"),
]

for i, (nome, status, itens) in enumerate(modulos):
    x = 0.4 + i * 4.3
    add_rect(slide, x, 1.4, 4.0, 5.5, BRANCO)
    add_rect(slide, x, 1.4, 4.0, 0.7, AZUL)
    add_text(slide, nome, x + 0.15, 1.48, 2.5, 0.55, 16, BRANCO, bold=True)
    add_text(slide, status, x + 0.15, 2.25, 3.7, 0.45, 11, AZUL, bold=True)
    add_text(slide, itens, x + 0.15, 2.75, 3.7, 3.8, 12, CINZA_TEXTO)


# ──────────────────────────────────────────────
# SLIDE 5 — CASE DE PRODUÇÃO
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK)
add_rect(slide, 0, 0, 13.33, 1.1, AZUL)
add_text(slide, "Resultado real. Provedor em produção.", 0.5, 0.2, 12, 0.8, 22, BRANCO, bold=True)
add_text(slide, "Provedor regional com 30.000 clientes ativos.", 0.8, 1.3, 11, 0.6, 14, CIANO)

metricas = [
    ("1.000", "leads por mês"),
    ("400",   "vendas digitais/mês"),
    ("R$99,90", "ticket médio"),
    ("2",     "pessoas no time\n(eram 8)"),
    ("R$284.400", "economia/ano\nem pessoal"),
]

for i, (numero, label) in enumerate(metricas):
    x = 0.5 + i * 2.55
    add_rect(slide, x, 2.1, 2.3, 2.2, RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, numero, x + 0.1, 2.2, 2.1, 0.9, 26, CIANO, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, 3.1, 2.1, 0.9, 11, CINZA_CLARO, align=PP_ALIGN.CENTER)

add_text(slide,
    "8 pessoas atendiam 1.000 leads por mês. Hoje 2 pessoas dão conta do mesmo volume.\n"
    "As outras 2 viraram backoffice estratégico. Função que antes não existia.",
    0.8, 4.7, 11.5, 1.3, 13, CINZA_CLARO)


# ──────────────────────────────────────────────
# SLIDE 6 — COMO FUNCIONA
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, CINZA_CLARO)
add_rect(slide, 0, 0, 13.33, 1.1, DARK)
add_text(slide, "Do lead ao HubSoft em minutos", 0.5, 0.2, 12, 0.8, 22, BRANCO, bold=True)

etapas = [
    ("1", "Lead no\nWhatsApp"),
    ("2", "Bot\nqualifica"),
    ("3", "Coleta\ndocumentos"),
    ("4", "Gera\ncontrato"),
    ("5", "Ativa no\nHubSoft"),
]

for i, (num, label) in enumerate(etapas):
    x = 0.6 + i * 2.5
    add_rect(slide, x, 2.0, 2.0, 2.0, AZUL)
    add_text(slide, num, x, 2.1, 2.0, 0.7, 28, BRANCO, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, 2.85, 2.0, 0.9, 13, BRANCO, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, "→", x + 2.05, 2.4, 0.4, 0.6, 22, AZUL, bold=True)

add_text(slide, "Sem planilha. Sem retrabalho. Sem lead perdido.", 0.8, 4.5, 11.5, 0.7, 16, DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Tudo integrado nativamente ao HubSoft.", 0.8, 5.1, 11.5, 0.6, 13, CINZA_TEXTO, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SLIDE 7 — PRECIFICAÇÃO
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK)
add_rect(slide, 0, 0, 13.33, 1.1, AZUL)
add_text(slide, "Preço que cresce com o provedor", 0.5, 0.2, 12, 0.8, 22, BRANCO, bold=True)

planos = [
    ("Starter", "R$397/mês", "ISPs pequenos\nAté 2 usuários\nFuncionalidades essenciais\nSuporte standard"),
    ("Start",   "R$797/mês", "ISPs em crescimento\nAté 5 usuários\nFuncionalidades completas\nSuporte standard"),
    ("Pro",     "R$1.397/mês","ISPs estabelecidos\nUsuários ilimitados\nIA + CRM Kanban\nSuporte prioritário"),
]

for i, (nome, preco, desc) in enumerate(planos):
    x = 0.6 + i * 4.2
    cor = AZUL if i == 1 else RGBColor(0x1E, 0x29, 0x3B)
    add_rect(slide, x, 1.4, 3.8, 5.5, cor)
    add_text(slide, nome, x + 0.15, 1.55, 3.5, 0.6, 18, BRANCO, bold=True)
    add_text(slide, preco, x + 0.15, 2.2, 3.5, 0.7, 22, CIANO, bold=True)
    add_text(slide, "+ transacional por uso", x + 0.15, 2.85, 3.5, 0.45, 10, CINZA_CLARO)
    add_text(slide, desc, x + 0.15, 3.4, 3.5, 3.0, 12, CINZA_CLARO)

add_text(slide, "Comercial | Marketing | CS — cada módulo tem seu próprio plano.", 0.8, 7.1, 11.5, 0.5, 11, CINZA_TEXTO, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SLIDE 8 — DIFERENCIAIS
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, CINZA_CLARO)
add_rect(slide, 0, 0, 13.33, 1.1, DARK)
add_text(slide, "Por que AuroraISP", 0.5, 0.2, 12, 0.8, 22, BRANCO, bold=True)

diferenciais = [
    ("HubSoft nativo",        "Integração profunda com o ERP mais usado por ISPs regionais. Nenhum concorrente faz isso."),
    ("Case real",             "Resultado comprovado em produção. Não é promessa. São números reais de um provedor com 30.000 clientes."),
    ("IA no processo",        "Qualificação, validação, sugestão de campanha e prevenção de churn com inteligência artificial."),
    ("Preço que faz sentido", "ROI positivo no primeiro mês. O transacional alinha o sucesso da Aurora com o sucesso do provedor."),
    ("Trial para HubSoft",    "14 dias gratuitos para quem já usa HubSoft. Em 14 dias dá para ver resultado real."),
    ("Hub completo",          "Comercial, Marketing e CS em uma plataforma. Do lead ao cliente fidelizado."),
]

for i, (titulo, desc) in enumerate(diferenciais):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.4 + row * 1.9
    add_rect(slide, x, y, 6.0, 1.65, BRANCO)
    add_rect(slide, x, y, 0.08, 1.65, AZUL)
    add_text(slide, titulo, x + 0.25, y + 0.1, 5.5, 0.5, 13, DARK, bold=True)
    add_text(slide, desc, x + 0.25, y + 0.6, 5.5, 0.9, 11, CINZA_TEXTO)


# ──────────────────────────────────────────────
# SLIDE 9 — CTA
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK)
add_rect(slide, 0, 0, 13.33, 1.1, AZUL)
add_text(slide, "Próximo passo", 0.5, 0.2, 12, 0.8, 22, BRANCO, bold=True)

add_text(slide, "São 30 minutos para você ver funcionando.", 0.8, 1.5, 11.5, 0.9, 26, BRANCO, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Demo ao vivo. Sem compromisso.", 0.8, 2.5, 11.5, 0.7, 18, CIANO, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 3.5, 6.3, 0.9, AZUL)
add_text(slide, "Agendar demonstração", 3.5, 3.58, 6.3, 0.75, 16, BRANCO, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Provedores HubSoft: 14 dias gratuitos para testar.", 0.8, 4.8, 11.5, 0.6, 13, CINZA_CLARO, align=PP_ALIGN.CENTER)
add_text(slide, "auroraisp.com.br", 0.8, 5.5, 11.5, 0.6, 14, CIANO, bold=True, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SALVAR
# ──────────────────────────────────────────────
output_dir = os.path.join(os.path.dirname(__file__), "..", "exports")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "auroraisp_pitch_deck.pptx")
prs.save(output_path)
print(f"Arquivo salvo em: {os.path.abspath(output_path)}")
